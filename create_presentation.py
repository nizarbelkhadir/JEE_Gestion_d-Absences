import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("La bibliothèque python-pptx n'est pas installée.")
    print("Veuillez l'installer en exécutant : pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    # Définition du format 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Couleurs du Design System (Glassmorphism Sombre)
    BG_COLOR = RGBColor(11, 15, 25)        # #0B0F19 - Bleu-Noir très sombre
    TEXT_WHITE = RGBColor(248, 250, 252)   # #F8FAFC - Blanc/Gris très clair
    ACCENT_PURPLE = RGBColor(139, 92, 246) # #8B5CF6 - Violet Électrique
    ACCENT_BLUE = RGBColor(59, 130, 246)   # #3B82F6 - Bleu Vif
    ACCENT_GREEN = RGBColor(16, 185, 129)  # #10B981 - Vert Succès
    ACCENT_RED = RGBColor(239, 68, 68)     # #EF4444 - Rouge Danger
    TEXT_GRAY = RGBColor(156, 163, 175)    # #9CA3AF - Gris Moyen
    CARD_BG = RGBColor(30, 41, 59)         # #1E293B - Gris-Bleu ardoise pour les cartes

    def add_custom_slide(title_text):
        # Utilisation d'un layout vierge (index 6)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Fond sombre uni
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Titre de la diapositive
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        return slide

    def add_card(slide, left, top, width, height, title, body_bullets, accent_color=ACCENT_BLUE):
        # Dessin du rectangle avec coins arrondis pour l'effet "carte"
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.fill.transparency = 0.45  # Transparence pour simuler le Glassmorphism
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)
        
        # Zone de texte à l'intérieur du conteneur
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)
        
        # Titre de la carte
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color
        p_title.space_after = Pt(12)
        
        # Puces textuelles courtes (pas de longs paragraphes)
        for bullet in body_bullets:
            p = tf.add_paragraph()
            p.text = "•  " + bullet
            p.font.name = 'Segoe UI'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1 : PAGE DE GARDE
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = BG_COLOR

    # Bandeau vertical décoratif lumineux à gauche
    glow = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_PURPLE
    glow.line.fill.background()

    # Titre & Sous-titre
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PLATEFORME DE GESTION DES ABSENCES"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)

    p_sub = tf.add_paragraph()
    p_sub.text = "Plateforme Multi-Rôles en Java EE 8 pur & JPA / Hibernate"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ACCENT_PURPLE

    # Informations du projet (Membres mis à jour)
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True

    p_team = tf_info.paragraphs[0]
    p_team.text = "Membres de l'équipe : Chafiq Mwiyeh, Adam Asaas, Nizar Belkhadir, Adam Bibilou"
    p_team.font.name = 'Segoe UI'
    p_team.font.size = Pt(16)
    p_team.font.bold = True
    p_team.font.color.rgb = TEXT_WHITE
    p_team.space_after = Pt(8)

    p_school = tf_info.add_paragraph()
    p_school.text = "Filière : 4IIR — EMSI  |  Encadrant : Pr. Houssam BAZZA"
    p_school.font.name = 'Segoe UI'
    p_school.font.size = Pt(16)
    p_school.font.color.rgb = ACCENT_BLUE
    p_school.space_after = Pt(8)

    p_date = tf_info.add_paragraph()
    p_date.text = "Soutenance de Projet de Fin d'Année  —  Juin 2026"
    p_date.font.name = 'Segoe UI'
    p_date.font.size = Pt(14)
    p_date.font.color.rgb = TEXT_GRAY

    # Notes de l'orateur (Chafiq)
    slide1.notes_slide.notes_text_frame.text = (
        "[Chafiq] Bonjour Monsieur Bazza, bonjour à tous.\n"
        "Nous sommes ravis de vous présenter aujourd'hui notre plateforme de gestion des absences multi-rôles, "
        "développée en Java EE 8 pur et JPA. Notre groupe est composé de Chafiq Mwiyeh, Adam Asaas, "
        "Nizar Belkhadir et Adam Bibilou.\n"
        "Je laisse maintenant la parole à Adam Bibilou pour vous présenter la problématique et nos choix technologiques."
    )

    # ----------------------------------------------------
    # SLIDE 2 : INTRODUCTION & PROBLEMATIQUE
    # ----------------------------------------------------
    slide2 = add_custom_slide("Introduction & Contexte")
    add_card(
        slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8),
        "Problématique de l'Assiduité",
        [
            "Suivi de l'assiduité universitaire manuel et fastidieux.",
            "Risque élevé de perte des justificatifs médicaux papier.",
            "Manque de réactivité (les étudiants sont avertis trop tard).",
            "Saisie chronophage pour les professeurs et l'administration."
        ],
        accent_color=ACCENT_RED
    )
    add_card(
        slide2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8),
        "Solution & Choix Techniques",
        [
            "Plateforme web interactive connectant les 3 rôles clés.",
            "Java EE 8 natif (Servlets, JSP, JSTL) sans framework Spring Boot.",
            "Persistance de données standardisée via JPA et Hibernate.",
            "Conteneurisation de la BDD PostgreSQL 15 sous Docker Compose."
        ],
        accent_color=ACCENT_GREEN
    )

    # Notes de l'orateur (Adam Bibilou)
    slide2.notes_slide.notes_text_frame.text = (
        "[Adam Bibilou] Le suivi de l'assiduité est un enjeu majeur dans notre établissement. "
        "Les processus manuels actuels engendrent des pertes de données, des retards de saisie et une charge administrative lourde. "
        "Notre solution est une plateforme web interactive qui connecte directement les étudiants, les professeurs et l'administration scolaire. "
        "Pour maîtriser les rouages des architectures logicielles, nous l'avons conçue en Java EE natif sans Spring Boot.\n"
        "Je laisse notre Scrum Master, Nizar, vous présenter comment nous avons organisé notre travail."
    )

    # ----------------------------------------------------
    # SLIDE 3 : METHODOLOGIE Agile (Scrum)
    # ----------------------------------------------------
    slide3 = add_custom_slide("Gestion de Projet Agile (Scrum)")
    add_card(
        slide3, Inches(0.6), Inches(1.6), Inches(3.8), Inches(4.8),
        "Sprint 1 : Socle & Design",
        [
            "Analyse fonctionnelle & diagrammes UML.",
            "Configuration PostgreSQL 15 / Docker Compose.",
            "Création des entités persistantes JPA.",
            "Mise en place du Design System CSS Glassmorphism."
        ],
        accent_color=ACCENT_BLUE
    )
    add_card(
        slide3, Inches(4.75), Inches(1.6), Inches(3.8), Inches(4.8),
        "Sprint 2 : Cœur Applicatif",
        [
            "Servlet unique de routage (Front Controller).",
            "Gestion des sessions et cookies Remember Me.",
            "Feuille d'appel interactive Enseignant.",
            "Module d'exportation CSV (BOM UTF-8 pour Excel)."
        ],
        accent_color=ACCENT_PURPLE
    )
    add_card(
        slide3, Inches(8.9), Inches(1.6), Inches(3.8), Inches(4.8),
        "Sprint 3 : Modération & Tests",
        [
            "Tableau de bord Étudiant réactif.",
            "Workflow de téléversement de justificatifs.",
            "Workflow de réclamations étudiants.",
            "Tests unitaires automatisés sous JUnit 4."
        ],
        accent_color=ACCENT_GREEN
    )

    # Notes de l'orateur (Nizar)
    slide3.notes_slide.notes_text_frame.text = (
        "[Nizar] En tant que Scrum Master, j'ai piloté l'équipe. Pour ce projet d'une durée de 3 semaines, "
        "nous avons adopté la méthodologie Scrum, découpée en 3 sprints d'une semaine.\n"
        "Le Sprint 1 a posé le socle de l'application (modèle BDD, JPA, Docker, CSS).\n"
        "Le Sprint 2 a implémenté le Front Controller et la feuille d'appel du professeur.\n"
        "Le Sprint 3 a achevé les dashboards étudiants, les justificatifs et les tests automatisés.\n"
        "Nos Daily Stand-ups quotidiens ont permis d'intégrer notre code sur Git sans aucun blocage.\n"
        "Nous allons maintenant vous présenter la phase de conception UML de nos modules."
    )

    # ----------------------------------------------------
    # SLIDE 4 : CONCEPTION & MODELISATION UML
    # ----------------------------------------------------
    slide4 = add_custom_slide("Conception & Modélisation UML")
    
    tx = slide4.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.133), Inches(0.6))
    p_uml = tx.text_frame.paragraphs[0]
    p_uml.text = "Modélisation des acteurs, du domaine de persistance et de la cinématique des requêtes :"
    p_uml.font.name = 'Segoe UI'
    p_uml.font.size = Pt(14)
    p_uml.font.color.rgb = TEXT_WHITE

    col_width = Inches(3.8)
    col_height = Inches(4.2)
    top_pos = Inches(1.7)

    diagrams = [
        ("use_case.png", Inches(0.6), "Diagramme de Cas d'Utilisation\n(use_case.png introuvable)"),
        ("class_diagram.png", Inches(4.75), "Diagramme de Classes JPA\n(class_diagram.png introuvable)"),
        ("sequence_diagram.png", Inches(8.9), "Diagramme de Séquence\n(sequence_diagram.png introuvable)")
    ]

    for img_name, left_pos, placeholder_text in diagrams:
        if os.path.exists(img_name):
            slide4.shapes.add_picture(img_name, left_pos, top_pos, width=col_width, height=col_height)
        else:
            sh = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, col_width, col_height)
            sh.fill.solid()
            sh.fill.fore_color.rgb = CARD_BG
            sh.line.color.rgb = ACCENT_PURPLE
            sh.line.width = Pt(1.5)
            
            tf_sh = sh.text_frame
            tf_sh.word_wrap = True
            p_sh = tf_sh.paragraphs[0]
            p_sh.alignment = PP_ALIGN.CENTER
            p_sh.text = placeholder_text
            p_sh.font.name = 'Segoe UI'
            p_sh.font.size = Pt(14)
            p_sh.font.color.rgb = TEXT_GRAY

    # Cadres descriptifs au bas (Noms mis à jour)
    labels = [
        ("1. Cas d'Utilisation (Chafiq)", Inches(0.6), ACCENT_BLUE),
        ("2. Classes JPA (Adam Asaas)", Inches(4.75), ACCENT_PURPLE),
        ("3. Séquence Requêtes (Nizar)", Inches(8.9), ACCENT_GREEN)
    ]
    for text_lbl, left_pos, color_lbl in labels:
        tx_lbl = slide4.shapes.add_textbox(left_pos, Inches(6.0), col_width, Inches(0.6))
        p_lbl = tx_lbl.text_frame.paragraphs[0]
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.text = text_lbl
        p_lbl.font.name = 'Segoe UI'
        p_lbl.font.size = Pt(14)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = color_lbl

    # Notes de l'orateur (Chafiq, Adam Asaas, Nizar)
    slide4.notes_slide.notes_text_frame.text = (
        "[Chafiq] J'ai conçu le Diagramme de Cas d'Utilisation. Il montre nos trois espaces sécurisés : "
        "l'administrateur gère le registre global et modère, l'enseignant fait l'appel, et l'étudiant justifie ses absences.\n\n"
        "[Adam Asaas] J'ai réalisé le Diagramme de Classes JPA. Il définit nos entités d'accès aux données : "
        "l'Etudiant possède des Absences, rattachées à une Seance de Cours d'un Professeur. C'est la base de notre mapping JPA/Hibernate.\n\n"
        "[Nizar] J'ai modélisé le Diagramme de Séquence. Il trace le flux complet d'une requête HTTP lors d'un upload de justificatif, "
        "depuis le clic sur le JSP, l'appel multipart intercepté par le contrôleur servlet, jusqu'à la persistance en base PostgreSQL via le DAO."
    )

    # ----------------------------------------------------
    # SLIDE 5 : ARCHITECTURE & DESIGN PATTERNS (PARTIE 1)
    # ----------------------------------------------------
    slide5 = add_custom_slide("Architecture & Design Patterns (1/2)")
    add_card(
        slide5, Inches(0.6), Inches(1.6), Inches(3.8), Inches(4.8),
        "1. Modèle-Vue-Contrôleur",
        [
            "Modèle : Entités JPA annotées gérant l'état.",
            "Vue : Pages JSPs dans WEB-INF/ pour le rendu graphique.",
            "Contrôleur : Servlet centrale orchestrent les flux.",
            "Séparation stricte assurant une maintenance facile."
        ],
        accent_color=ACCENT_BLUE
    )
    add_card(
        slide5, Inches(4.75), Inches(1.6), Inches(3.8), Inches(4.8),
        "2. Front Controller",
        [
            "Servlet unique (AbsenceServlet) interceptant les URL /absences.",
            "Vérification centralisée des sessions utilisateurs.",
            "Aiguillage réactif basé sur le paramètre d'action.",
            "Sécurité renforcée (Session Gating par rôle)."
        ],
        accent_color=ACCENT_PURPLE
    )
    add_card(
        slide5, Inches(8.9), Inches(1.6), Inches(3.8), Inches(4.8),
        "3. Data Access Object (DAO)",
        [
            "DAO isole les requêtes JPA (JPQL Hibernate).",
            "Aucun code de persistance dans le contrôleur.",
            "Gestion explicite des transactions en local.",
            "Indépendance complète vis-à-vis du stockage physique."
        ],
        accent_color=ACCENT_GREEN
    )

    # Notes de l'orateur (Chafiq et Nizar)
    slide5.notes_slide.notes_text_frame.text = (
        "[Chafiq] Pour notre architecture, nous avons implémenté le pattern MVC. "
        "Les entités persistantes JPA représentent notre Modèle. Les pages JSP sous WEB-INF constituent la Vue. "
        "La servlet unique joue le rôle de Contrôleur. Cette séparation évite de mêler logique métier et affichage.\n\n"
        "[Nizar] Pour centraliser le routage et sécuriser les accès, j'ai mis en place un Front Controller dans notre servlet unique. "
        "Elle vérifie la validité des sessions et redirige vers le login si besoin. Nous utilisons également le pattern DAO "
        "pour encapsuler nos requêtes JPA et nos transactions Hibernate sans polluer le contrôleur."
    )

    # ----------------------------------------------------
    # SLIDE 6 : PERSISTANCE & GESTION DES THREADS (PARTIE 2)
    # ----------------------------------------------------
    slide6 = add_custom_slide("Architecture & Design Patterns (2/2)")
    add_card(
        slide6, Inches(0.6), Inches(1.6), Inches(3.8), Inches(4.8),
        "1. Singleton (JPAUtil)",
        [
            "EntityManagerFactory est coûteuse à instancier.",
            "Instanciée une seule fois au démarrage par ServletContextListener.",
            "Partagée sous forme d'instance unique globale.",
            "Économise les ressources mémoire du serveur Tomcat."
        ],
        accent_color=ACCENT_BLUE
    )
    add_card(
        slide6, Inches(4.75), Inches(1.6), Inches(3.8), Inches(4.8),
        "2. ThreadLocal Context",
        [
            "L'EntityManager standard n'est pas Thread-Safe.",
            "Tomcat traite chaque client dans un Thread distinct.",
            "ThreadLocal isole l'EntityManager par thread de requête.",
            "Évite les corruptions et conflits d'accès concurrents."
        ],
        accent_color=ACCENT_PURPLE
    )
    add_card(
        slide6, Inches(8.9), Inches(1.6), Inches(3.8), Inches(4.8),
        "3. Sessions & Cookies",
        [
            "HTTP est sans état (stateless).",
            "Maintien de session via HttpSession (JSESSIONID).",
            "Cookies sécurisés 'Remember Me' persistants.",
            "Remplissage automatique du formulaire de login."
        ],
        accent_color=ACCENT_GREEN
    )

    # Notes de l'orateur (Adam Asaas & Adam Bibilou)
    slide6.notes_slide.notes_text_frame.text = (
        "[Adam Asaas] Pour optimiser les ressources de notre serveur, nous utilisons le pattern Singleton dans notre classe JPAUtil. "
        "L'EntityManagerFactory n'est instancié qu'une seule fois au chargement du contexte web. "
        "Pour maintenir l'authentification active entre les pages, nous utilisons la session HTTP ainsi que des Cookies 'Remember Me' persistants.\n\n"
        "[Adam Bibilou] L'EntityManager standard de JPA n'étant pas thread-safe, j'ai configuré le pattern ThreadLocal. "
        "Sous Tomcat, chaque client exécute ses requêtes sur un thread propre. Le ThreadLocal lie un EntityManager au thread courant. "
        "Il est ouvert au début de la requête et fermé à la fin, garantissant une isolation complète de nos transactions de données."
    )

    # ----------------------------------------------------
    # SLIDE 7 : ROBUSTESSE & TESTS UNITAIRES (JUNIT)
    # ----------------------------------------------------
    slide7 = add_custom_slide("Robustesse & Tests Unitaires (JUnit)")
    add_card(
        slide7, Inches(0.6), Inches(1.6), Inches(3.8), Inches(4.8),
        "Cadre de Validation",
        [
            "Suite de tests unitaires automatisés sous JUnit 4.",
            "Ciblent le comportement logique des modèles et utilitaires.",
            "Exécutés en isolation complète avant chaque build.",
            "Garantissent la non-régression lors des refactorings."
        ],
        accent_color=ACCENT_BLUE
    )
    add_card(
        slide7, Inches(4.75), Inches(1.6), Inches(3.8), Inches(4.8),
        "Modèles Validés",
        [
            "EtudiantTest : valide les constructeurs, accesseurs et la capitalisation automatique du nom de famille.",
            "SeanceTest : vérifie le formateur de date (dd/MM/yyyy) et la génération du libellé séance."
        ],
        accent_color=ACCENT_PURPLE
    )
    add_card(
        slide7, Inches(8.9), Inches(1.6), Inches(3.8), Inches(4.8),
        "Persistance Validée",
        [
            "JPAUtilTest : valide la logique de persistance.",
            "Vérifie la levée d'une IllegalStateException si la fabrique n'est pas démarrée.",
            "Assure un démarrage sain du serveur."
        ],
        accent_color=ACCENT_GREEN
    )

    # Notes de l'orateur (Nizar)
    slide7.notes_slide.notes_text_frame.text = (
        "[Nizar] Pour garantir la qualité et la robustesse de notre code, j'ai écrit une suite de tests unitaires avec JUnit 4. "
        "Nous testons en isolation nos classes métiers et utilitaires.\n"
        "EtudiantTest vérifie la mise en majuscule automatique du nom lors de l'appel à getNomComplet().\n"
        "SeanceTest valide le formateur de date personnalisé en dd/MM/yyyy.\n"
        "JPAUtilTest s'assure du bon comportement de notre fabrique en cas de défaut d'initialisation.\n"
        "Cela protège notre application de toute régression technique."
    )

    # ----------------------------------------------------
    # SLIDE 8 : DEMONSTRATION & EXECUTION (VIDEO)
    # ----------------------------------------------------
    slide8 = add_custom_slide("Démonstration & Exécution")
    
    video_path = "app.mp4"
    v_left = Inches(1.66)
    v_top = Inches(1.4)
    v_width = Inches(10.0)
    v_height = Inches(5.3)

    if os.path.exists(video_path):
        try:
            slide8.shapes.add_movie(
                video_path,
                v_left, v_top, v_width, v_height,
                mime_type='video/mp4'
            )
            tx_v = slide8.shapes.add_textbox(Inches(1.66), Inches(6.75), Inches(10.0), Inches(0.5))
            p_v = tx_v.text_frame.paragraphs[0]
            p_v.alignment = PP_ALIGN.CENTER
            p_v.text = "Note : Pour lire automatiquement en boucle, réglez l'onglet Lecture -> Démarrer : Automatiquement."
            p_v.font.name = 'Segoe UI'
            p_v.font.size = Pt(11)
            p_v.font.color.rgb = TEXT_GRAY
        except Exception as e:
            sh_v = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, v_left, v_top, v_width, v_height)
            sh_v.fill.solid()
            sh_v.fill.fore_color.rgb = CARD_BG
            sh_v.line.color.rgb = ACCENT_RED
            sh_v.text = f"Erreur lors de l'insertion de la vidéo :\n{str(e)}"
    else:
        placeholder = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, v_left, v_top, v_width, v_height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = CARD_BG
        placeholder.line.color.rgb = ACCENT_PURPLE
        placeholder.line.width = Pt(2)
        
        tf_v = placeholder.text_frame
        tf_v.word_wrap = True
        p_v1 = tf_v.paragraphs[0]
        p_v1.alignment = PP_ALIGN.CENTER
        p_v1.text = "EMPLACEMENT DE LA VIDÉO DE DÉMO (app.mp4)"
        p_v1.font.name = 'Segoe UI'
        p_v1.font.size = Pt(20)
        p_v1.font.bold = True
        p_v1.font.color.rgb = ACCENT_PURPLE
        p_v1.space_after = Pt(20)
        
        p_v2 = tf_v.add_paragraph()
        p_v2.alignment = PP_ALIGN.CENTER
        p_v2.text = (
            "Pour intégrer votre vidéo :\n"
            "1. Enregistrez votre vidéo de démo sous le nom 'app.mp4' dans ce dossier.\n"
            "2. Relancez ce script de génération.\n\n"
            "Conseil PowerPoint : Allez dans l'onglet 'Lecture' de la vidéo, puis réglez\n"
            "le démarrage sur 'Automatiquement' et cochez 'Boucler jusqu'à l'arrêt' et sans son."
        )
        p_v2.font.name = 'Segoe UI'
        p_v2.font.size = Pt(14)
        p_v2.font.color.rgb = TEXT_WHITE

    # Notes de l'orateur partagées (Chafiq, Nizar, Adam Asaas, Adam Bibilou)
    slide8.notes_slide.notes_text_frame.text = (
        "[Chafiq] À l'écran, nous voyons la connexion. Je m'authentifie en Administrateur pour consulter le registre global des absences déclarées. "
        "Je réalise ensuite l'exportation au format CSV. Le fichier généré est directement exploitable sous Microsoft Excel grâce au BOM UTF-8.\n\n"
        "[Nizar] Nous passons ensuite sur l'espace Enseignant. Le professeur choisit sa séance dans sa liste personnalisée et remplit sa feuille d'appel "
        "interactive en cochant les étudiants absents. La validation soumet les absences qui sont persistées en bloc en base de données.\n\n"
        "[Adam Asaas] Nous arrivons sur le tableau de bord de l'Étudiant. Il peut consulter son taux d'assiduité en temps réel. "
        "L'étudiant choisit une absence injustifiée et téléverse sa pièce jointe (PDF). La demande de justification passe alors au statut 'En attente'.\n\n"
        "[Adam Bibilou] Enfin, l'étudiant peut soumettre une réclamation en cas d'erreur de saisie du professeur. "
        "Sur l'espace d'administration, on voit le portail de modération : l'administrateur valide le justificatif médical ou accepte la réclamation, "
        "ce qui met à jour le statut ou supprime physiquement l'absence de la base PostgreSQL."
    )

    # ----------------------------------------------------
    # SLIDE 9 : CONCLUSION & REMERCIEMENTS
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.rgb = BG_COLOR

    # Bandeau vertical décoratif lumineux à gauche
    glow9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    glow9.fill.solid()
    glow9.fill.fore_color.rgb = ACCENT_GREEN
    glow9.line.fill.background()

    # Titre de conclusion
    c_box = slide9.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(2.5))
    tf_c = c_box.text_frame
    tf_c.word_wrap = True
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "CONCLUSION & PERSPECTIVES"
    p_c1.font.name = 'Segoe UI'
    p_c1.font.size = Pt(36)
    p_c1.font.bold = True
    p_c1.font.color.rgb = ACCENT_GREEN
    p_c1.space_after = Pt(14)

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "•  Validation réussie de l'architecture JEE 8 pur et de la persistance JPA isolée."
    p_c2.font.name = 'Segoe UI'
    p_c2.font.size = Pt(16)
    p_c2.font.color.rgb = TEXT_WHITE
    p_c2.space_after = Pt(8)

    p_c3 = tf_c.add_paragraph()
    p_c3.text = "•  Perspectives : Migration vers une architecture Microservices et intégration de notifications Push/SMS."
    p_c3.font.name = 'Segoe UI'
    p_c3.font.size = Pt(16)
    p_c3.font.color.rgb = TEXT_WHITE

    # Message de remerciement
    t_box = slide9.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.5))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "Merci pour votre attention."
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_after = Pt(8)

    p_q = tf_t.add_paragraph()
    p_q.text = "Nous sommes ouverts à vos questions."
    p_q.font.name = 'Segoe UI'
    p_q.font.size = Pt(20)
    p_q.font.color.rgb = ACCENT_GREEN

    # Notes de l'orateur (Adam Bibilou)
    slide9.notes_slide.notes_text_frame.text = (
        "[Adam Bibilou] Pour conclure, ce projet nous a permis de mettre en pratique et d'apprécier la puissance des standards natifs JEE. "
        "L'application est pleinement opérationnelle.\n"
        "Dans le futur, nous aimerions faire évoluer cette plateforme vers une architecture microservices "
        "et intégrer un système de notifications instantanées par e-mail ou SMS pour alerter les étudiants.\n"
        "Nous vous remercions pour votre attention et sommes maintenant ouverts à vos questions."
    )

    # ----------------------------------------------------
    # SAUVEGARDE DE LA PRESENTATION
    # ----------------------------------------------------
    output_filename = "presentation_soutenance.pptx"
    prs.save(output_filename)
    print(f"Présentation générée avec succès : {output_filename}")

if __name__ == "__main__":
    create_presentation()
